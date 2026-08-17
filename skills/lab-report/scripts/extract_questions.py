"""Dump the text of a question paper (PDF / PPTX / DOCX / TXT) and flag the
questions whose answer depends on the student.

    python extract_questions.py "4th Sem/OS/Lab Works1.pdf"
    python extract_questions.py "DBMS Lab 1.pptx" --raw     # text only, no marker report

Output:
  ==== page/slide N ====      the text, in reading order
  ==== STUDENT-SPECIFIC ====  every line that mentions the student, with its line number
  ==== IMAGES ====            pages that carry pictures - read those pages with the Read
                              tool (Read supports PDFs and images) because the text dump
                              silently loses diagrams, family trees, DFA drawings, etc.
"""
from __future__ import annotations

import argparse
import re
import sys
import zipfile
from pathlib import Path

MARKERS = [
    r"your\s*name", r"your\s*roll", r"yourname", r"<\s*your", r"your\s*full\s*name",
    r"first\s*letter", r"initial", r"your\s*choice", r"your\s*section", r"your\s*id",
    r"roll\s*number", r"roll\s*no", r"folder\s*name", r"file\s*path", r"your\s*own",
    r"at\s*the\s*end\s*in\s*each\s*output", r"print\s+your", r"lab\s*no\.?\s*:\s*name",
]
MARKER_RE = re.compile("|".join(MARKERS), re.IGNORECASE)


def from_pdf(path: Path) -> tuple[list[tuple[str, str]], list[int]]:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        import subprocess
        text = subprocess.run(["pdftotext", "-layout", str(path), "-"],
                              capture_output=True, text=True, encoding="utf-8", errors="replace").stdout
        return [("whole file", text)], []
    doc = fitz.open(path)
    pages, with_images = [], []
    for number, page in enumerate(doc, 1):
        pages.append((f"page {number}", page.get_text()))
        if page.get_images(full=True):
            with_images.append(number)
    return pages, with_images


def from_pptx(path: Path) -> tuple[list[tuple[str, str]], list[int]]:
    from pptx import Presentation
    presentation = Presentation(path)
    slides, with_images = [], []
    for number, slide in enumerate(presentation.slides, 1):
        chunks = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            chunks.append("[notes] " + slide.notes_slide.notes_text_frame.text)
        slides.append((f"slide {number}", "\n".join(chunks)))
        if any(shape.shape_type == 13 for shape in slide.shapes):     # 13 = PICTURE
            with_images.append(number)
    return slides, with_images


def from_docx(path: Path) -> tuple[list[tuple[str, str]], list[int]]:
    try:
        import docx
        return [("whole file", "\n".join(p.text for p in docx.Document(path).paragraphs))], []
    except ImportError:
        with zipfile.ZipFile(path) as archive:
            xml = archive.read("word/document.xml").decode("utf-8", errors="replace")
        xml = re.sub(r"</w:p>", "\n", xml)
        return [("whole file", re.sub(r"<[^>]+>", "", xml))], []


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("path")
    parser.add_argument("--raw", action="store_true")
    args = parser.parse_args()
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

    path = Path(args.path).resolve()
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        sections, with_images = from_pdf(path)
    elif suffix == ".pptx":
        sections, with_images = from_pptx(path)
    elif suffix in (".docx",):
        sections, with_images = from_docx(path)
    elif suffix in (".txt", ".md", ".sql"):
        sections, with_images = [("whole file", path.read_text(encoding="utf-8", errors="replace"))], []
    else:
        sys.exit(f"Unsupported {suffix}. Images and scans: use the Read tool on the file directly.")

    lines: list[str] = []
    for label, text in sections:
        print(f"\n==== {label} ====")
        for line in text.splitlines():
            print(line)
            lines.append(line)

    if args.raw:
        return

    hits = [(number, line.strip()) for number, line in enumerate(lines, 1)
            if line.strip() and MARKER_RE.search(line)]
    print("\n==== STUDENT-SPECIFIC ====")
    if hits:
        for number, line in hits:
            print(f"L{number}: {line}")
    else:
        print("none found - still read the instructions block yourself before deciding")

    print("\n==== IMAGES ====")
    print(", ".join(str(n) for n in with_images) if with_images
          else "no embedded images detected")


if __name__ == "__main__":
    main()
